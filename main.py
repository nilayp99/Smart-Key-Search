import sys

import re

import fitz

from PyQt5.QtWidgets import (

    QApplication, QMainWindow, QVBoxLayout, QWidget, QLabel, QFileDialog, QGraphicsTransform, QSizePolicy, QProgressBar,

    QPushButton, QTextEdit, QScrollArea, QListWidget, QListWidgetItem, QHBoxLayout, QSplitter, QItemDelegate, QGroupBox

)

from PyQt5.QtGui import QPixmap, QImage, QColor, QFont, QTransform, QPainter, QIcon

from PyQt5.QtCore import Qt, QSize, QThread, pyqtSignal

import os

import io

from reportlab.lib.pagesizes import letter

from reportlab.pdfgen import canvas

from reportlab.lib import colors

from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

from reportlab.lib.colors import yellow, cyan, magenta, red, blue, pink, orange, green

from docx import Document

from pptx import Presentation

import extract_msg

from python_calamine import CalamineWorkbook

 

 

class HighlightThread(QThread):

    highlight_complete = pyqtSignal(str)  # Emits the path of the highlighted file when done

    error = pyqtSignal(str)

 

    def __init__(self, file_path, keywords, regex_string, file_extension, parent=None):

        super().__init__(parent)

        self.file_path = file_path

        self.keywords = keywords

        self.regex_string = regex_string

        self.file_extension = file_extension

 

    def run(self):

        try:

            if self.file_extension in ['.pdf', '.txt']:

                highlighted_file_path = self.highlight_keywords_in_pdf(self.file_path, self.keywords)

                self.highlight_complete.emit(highlighted_file_path)

            elif self.file_extension == '.docx':

                highlighted_file_path = self.highlight_keywords_in_docx(self.file_path, self.keywords)

                self.highlight_complete.emit(highlighted_file_path)

            elif self.file_extension == '.pptx':

                highlighted_file_path = self.highlight_keywords_in_pptx(self.file_path, self.keywords)

                self.highlight_complete.emit(highlighted_file_path)

            elif self.file_extension == '.xlsx':

                highlighted_file_path = self.highlight_keywords_in_xlsx(self.file_path, self.keywords)

                self.highlight_complete.emit(highlighted_file_path)

            elif self.file_extension == '.msg':

                highlighted_file_path = self.highlight_keywords_in_msg(self.file_path, self.keywords)

                self.highlight_complete.emit(highlighted_file_path)

        except Exception as e:

            self.error.emit(str(e))

 

    def match_regex(self, word):

        try:

            # Compile the regex pattern to check if it is valid

            pattern = re.compile(self.regex_string)

            # Match the word against the compiled pattern

            return bool(pattern.fullmatch(word))

        except re.error as e:

            # Handle invalid regex patterns

            raise ValueError(f"Invalid regex pattern: {e}")

 

    def highlight_keywords_in_pdf(self, pdf_path, keywords):

        base, _ = os.path.splitext(pdf_path)

        highlighted_pdf_path = f"{base}_highlighted_{'_'.join(keywords)}.pdf"

 

        document = fitz.open(pdf_path)

        highlight_colors = [

                            fitz.utils.getColor("magenta"),

                            fitz.utils.getColor("red"),

                            fitz.utils.getColor("blue"),

                            fitz.utils.getColor("pink"),

                            fitz.utils.getColor("orange"),

                            fitz.utils.getColor("green"),

                            fitz.utils.getColor("yellow"),

                            fitz.utils.getColor("cyan")

                           ]

 

 

        for page_num in range(len(document)):

            page = document.load_page(page_num)

 

            for i, keyword in enumerate(keywords):

                if keyword != '' and keyword != '((regex))':

                    text_instances = page.search_for(keyword, quads = True)

 

                    for inst in text_instances:

                        highlight = page.add_highlight_annot(inst)

                        highlight.set_colors({"stroke": highlight_colors[i % len(highlight_colors)]})

                        highlight.update()

 

                # Highlight words matching the regex

                if keyword == '((regex))':

                    words = page.get_text("words")  # Extract all words on the page

                    for word_data in words:

                        if self.match_regex(word_data[4]):  # Match the word against the regex

                            inst = fitz.Rect(word_data[0], word_data[1], word_data[2], word_data[3])

                            highlight = page.add_highlight_annot(inst)

                            highlight.set_colors({"stroke": highlight_colors[len(keywords)-1 % len(highlight_colors)]})

                            highlight.update()

 

        document.save(highlighted_pdf_path)

        document.close()

 

        return highlighted_pdf_path

 

    def convert_to_column_alphabet(self, num):

        """Convert column number to Excel column alphabet."""

        result = ''

        while num > 0:

            num, remainder = divmod(num - 1, 26)

            result = chr(65 + remainder) + result

        return result

 

 

    def highlight_keywords_in_xlsx(self, xlsx_path, keywords):

        # Save the PDF to a file

        base, _ = os.path.splitext(xlsx_path)

        summary_pdf_path = f"{base}_highlighted_{'_'.join(keywords)}.pdf"

 

        wb = CalamineWorkbook.from_path(xlsx_path)

        keyword_positions = {keyword: [] for keyword in keywords}

 

        for sheet_name in wb.sheet_names:

            ws = wb.get_sheet_by_name(sheet_name)

 

            for i, row in enumerate(ws.to_python(), start=1):

                for j, cell in enumerate(row, start=1):

                    if cell:  # Check if cell has a value

                        cell_value = str(cell).lower()  # Convert cell value to lowercase

                        for keyword in keywords:

                            if keyword != '' and keyword != '((regex))':

                                if keyword.lower() in cell_value:  # Case-insensitive search

                                    # Add the keyword, sheet name, row number, and column number

                                    cell_position = {

                                        "sheet": sheet_name,

                                        "row": i,     # Row number from cell object

                                        "col": self.convert_to_column_alphabet(j),  # Column number from cell object

                                        "value": str(cell)

                                    }

                                    keyword_positions[keyword].append(cell_position)

                            if keyword == '((regex))':

                                for word in str(cell).split():

                                    if self.match_regex(word):

                                        cell_position = {

                                            "sheet": sheet_name,

                                            "row": i,     # Row number from cell object

                                            "col": self.convert_to_column_alphabet(j),  # Column number from cell object

                                            "value": str(cell)

                                        }

                                        keyword_positions[keyword].append(cell_position)

                                        break

 

 

        # Filter out empty keyword lists

        keyword_positions = {k: v for k, v in keyword_positions.items() if v}

 

        # Create an in-memory PDF buffer

        buffer = io.BytesIO()

        pdf = SimpleDocTemplate(buffer, pagesize=letter)

        styles = getSampleStyleSheet()

        story = []

 

        # Add a title to the PDF

        story.append(Paragraph("Keyword Search Results for XLSX", styles['Title']))

        story.append(Spacer(1, 12))

 

        # Iterate over the found keyword positions

        for keyword, positions in keyword_positions.items():

            story.append(Paragraph(f"Keyword: {keyword}", styles['Heading2']))

            story.append(Spacer(1, 12))

 

            # Create a table with headers for displaying the results

            data = [["Sheet Name", "Row", "Column", "Value"]]  # Header row

 

            for position in positions:

                sheet = position['sheet']

                row = position['row']

                col = position['col']

                value = position['value']

                data.append([sheet, row, col, value])  # Append the found position details

 

            # Wrap long text values using Paragraph and Table

            wrapped_data = []

            for row_data in data:

                wrapped_row = []

                for cell in row_data:

                    wrapped_row.append(Paragraph(str(cell), styles['Normal']))  # Convert cell to string for Paragraph

                wrapped_data.append(wrapped_row)

 

            # Create the table

            table = Table(wrapped_data, colWidths=[100, 50, 50, 300])  # Set column widths

            table.setStyle(TableStyle([

                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),

                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),

                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),

                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),

                ('FONTSIZE', (0, 0), (-1, 0), 12),

                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),

                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),

                ('GRID', (0, 0), (-1, -1), 1, colors.black),

                ('VALIGN', (0, 0), (-1, -1), 'TOP'),  # Vertically align text

            ]))

 

            # Add the table to the story

            story.append(table)

            story.append(Spacer(1, 24))

 

        # Build the PDF

        pdf.build(story)

 

        with open(summary_pdf_path, 'wb') as f:

            f.write(buffer.getvalue())

 

        return summary_pdf_path

 

    def highlight_keywords_in_pptx(self, pptx_path, keywords):

        # Save the highlighted PDF

        base, _ = os.path.splitext(pptx_path)

        highlighted_pdf_path = f"{base}_highlighted_{'_'.join(keywords)}.pdf"

 

        prs = Presentation(pptx_path)

 

        # Create an in-memory PDF buffer

        buffer = io.BytesIO()

        pdf = SimpleDocTemplate(buffer, pagesize=letter)

        styles = getSampleStyleSheet()

        story = []

        colors = ['magenta', 'red', 'blue', 'pink', 'orange', 'green','yellow', 'cyan']

 

        for slide in prs.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    words = shape.text.split()

                    highlighted_words = []

 

                    for word in words:

                        found_keyword = False

                        for i, keyword in enumerate(keywords):

                            if keyword != '' and keyword != '((regex))':

                                if keyword.lower() in word.lower():

                                    text_split = word.lower().split(keyword.lower())

                                    for j in range(len(text_split)):

                                        highlighted_words.append(text_split[j])

                                        if j != len(text_split)-1:

                                            highlighted_word = f'<span bgcolor="{colors[i % len(colors)]}">{keyword.lower()}</span>'

                                            highlighted_words.append(highlighted_word)

                                    found_keyword = True

                                    break

                            if keyword == '((regex))' and self.match_regex(word):

                                highlighted_word = f'<span bgcolor="{colors[len(keywords)-1 % len(colors)]}">{word}</span>'

                                highlighted_words.append(highlighted_word)

                                found_keyword = True

                                break

 

                        if not found_keyword:

                            highlighted_words.append(word)

 

                    highlighted_text = Paragraph(' '.join(highlighted_words), styles['Normal'])

                    story.append(highlighted_text)

                    story.append(Spacer(1, 12))

 

        pdf.build(story)

 

        with open(highlighted_pdf_path, 'wb') as f:

            f.write(buffer.getvalue())

 

        return highlighted_pdf_path

 

 

    def highlight_keywords_in_docx(self, docx_path, keywords):

        # Save the highlighted PDF

        base, _ = os.path.splitext(docx_path)

        highlighted_pdf_path = f"{base}_highlighted_{'_'.join(keywords)}.pdf"

 

        doc = Document(docx_path)

   

        # Create an in-memory PDF

        buffer = io.BytesIO()

        pdf = SimpleDocTemplate(buffer, pagesize=letter)

        styles = getSampleStyleSheet()

        story = []

        colors = ['magenta', 'red', 'blue', 'pink', 'orange', 'green', 'yellow', 'cyan']

 

        for paragraph in doc.paragraphs:

            words = paragraph.text.split()

            highlighted_words = []

 

            for word in words:

                found_keyword = False

                for i, keyword in enumerate(keywords):

                    if keyword != '' and keyword != '((regex))':

                        if keyword.lower() in word.lower():

                            text_split = word.lower().split(keyword.lower())

                            for j in range(len(text_split)):

                                highlighted_words.append(text_split[j])

                                if j != len(text_split)-1:

                                    highlighted_word = f'<span bgcolor="{colors[i % len(colors)]}">{keyword.lower()}</span>'

                                    highlighted_words.append(highlighted_word)

                            found_keyword = True

                            break

                    if keyword == '((regex))' and self.match_regex(word):

                        highlighted_word = f'<span bgcolor="{colors[len(keywords)-1 % len(colors)]}">{word}</span>'

                        highlighted_words.append(highlighted_word)

                        found_keyword = True

                        break

 

                if not found_keyword:

                    highlighted_words.append(word)

 

            highlighted_text = Paragraph(' '.join(highlighted_words), styles['Normal'])

            story.append(highlighted_text)

            story.append(Spacer(1, 12))

 

        pdf.build(story)

 

        with open(highlighted_pdf_path, 'wb') as f:

            f.write(buffer.getvalue())

 

        return highlighted_pdf_path

 

    def highlight_keywords_in_msg(self, msg_path, keywords):

        try:

            # Save the highlighted PDF

            base, _ = os.path.splitext(msg_path)

            highlighted_pdf_path = f"{base}_highlighted_{'_'.join(keywords)}.pdf"

 

            msg = extract_msg.Message(msg_path)

            msg_content = msg.body  # Extract the email content (body)

 

            # Create an in-memory PDF

            buffer = io.BytesIO()

            pdf = SimpleDocTemplate(buffer, pagesize=letter)

            styles = getSampleStyleSheet()

            story = []

            colors = ['magenta', 'red', 'blue', 'pink', 'orange', 'green', 'yellow', 'cyan']

 

            # Split the email content into words and highlight the keywords

            words = msg_content.split()

            highlighted_words = []

 

            for word in words:

                found_keyword = False

                for i, keyword in enumerate(keywords):

                    if keyword != '' and keyword != '((regex))':

                        if keyword.lower() in word.lower():

                            text_split = word.lower().split(keyword.lower())

                            for j in range(len(text_split)):

                                highlighted_words.append(text_split[j])

                                if j != len(text_split)-1:

                                    highlighted_word = f'<span bgcolor="{colors[i % len(colors)]}">{keyword.lower()}</span>'

                                    highlighted_words.append(highlighted_word)

                            found_keyword = True

                            break

                    if keyword == '((regex))' and self.match_regex(word):

                        highlighted_word = f'<span bgcolor="{colors[len(keywords)-1 % len(colors)]}">{word}</span>'

                        highlighted_words.append(highlighted_word)

                        found_keyword = True

                        break

 

                if not found_keyword:

                    highlighted_words.append(word)

 

            # Create a paragraph for the highlighted content

            highlighted_text = Paragraph(' '.join(highlighted_words), styles['Normal'])

            story.append(highlighted_text)

            story.append(Spacer(1, 12))

 

            # Build the PDF

            pdf.build(story)

 

            with open(highlighted_pdf_path, 'wb') as f:

 

                f.write(buffer.getvalue())

 

            return highlighted_pdf_path

 

        except Exception as e:

            return None

 

 

 

class SearchThread(QThread):

    progress_within_file = pyqtSignal(int)  # Progress within a single file

    update_current_file = pyqtSignal(str)   # Signal for updating current file

    file_found = pyqtSignal(str)            # Signal when a file with the keyword is found

    search_complete = pyqtSignal()          # Signal when the search is done

    files_searched = pyqtSignal(int, int)   # Signal for updating number of searched files

 

    def __init__(self, folder_path, keywords, regex_string, file_extensions, parent=None):

        super().__init__(parent)

        self.folder_path = folder_path

        self.keywords = keywords

        self.file_extensions = file_extensions

        self.regex_string = regex_string

 

    def run(self):

        files_to_search = []

 

        # Collect files to search, along with their sizes

        for root, _, files in os.walk(self.folder_path):

            for file in files:

                if any(file.lower().endswith(ext) for ext in self.file_extensions):

                    full_path = os.path.join(root, file)

                    file_size = os.path.getsize(full_path)  # Get file size

                    files_to_search.append((full_path, file_size))

 

        # Sort files by size (smallest first)

        files_to_search.sort(key=lambda x: x[1])

 

        total_files = len(files_to_search)

        if total_files == 0:

            self.search_complete.emit()

            return

 

        for idx, (file_path, _) in enumerate(files_to_search):

            self.update_current_file.emit(file_path)

            keyword_found = self.search_within_file(file_path, self.keywords, self.regex_string)

 

            if keyword_found:

                self.file_found.emit(file_path)

           

            # Emit the progress for total files searched

            self.files_searched.emit(idx + 1, total_files)

 

        self.search_complete.emit()

 

    def search_within_file(self, file_path, keywords, regex_string):

        file_name, file_extension = os.path.splitext(file_path)

        file_extension = file_extension.lower()

 

        found_keywords = {}

 

        # Initialize found_keywords dictionary to track each keyword

        if '' not in keywords:

            found_keywords = {keyword.lower(): False for keyword in keywords}

 

        if regex_string != '':

            try:

                pattern = re.compile(regex_string)

            except:

                return False

            found_keywords['((regex))'] = False

 

        if found_keywords:

            # For .xlsx files

            if file_extension == '.xlsx':

                try:

                    # Open the workbook using calamine

                    wb = CalamineWorkbook.from_path(file_path)

                    found_keywords = {keyword.lower(): False for keyword in keywords}

 

                    if regex_string != '':

                        found_keywords['((regex))'] = False

 

                    total_rows = 0

                    # Calculate total number of rows across all sheets

                    # Hard coded for faster processing

                    for sheet_name in wb.sheet_names:

                        sheet = wb.get_sheet_by_name(sheet_name)

                        total_rows += len(sheet.to_python())

 

                    processed_rows = 0

 

                    # Iterate through each sheet and each row

                    for sheet_name in wb.sheet_names:

                        sheet = wb.get_sheet_by_name(sheet_name)

 

                        for row in sheet.to_python():

                            for cell in row:

                                if cell:

                                    cell_value = str(cell).lower()

                                    for keyword in keywords:

                                        if keyword.lower() in cell_value:

                                            found_keywords[keyword.lower()] = True

                                        if keyword.lower() == '((regex))':

                                            for word in cell_value.split():

                                                # Match the word against the compiled pattern

                                                if bool(pattern.fullmatch(word)):

                                                    found_keywords['((regex))'] = True

                                                    break

 

 

                            # Emit progress after processing each row

                            processed_rows += 1

                            self.progress_within_file.emit(int((processed_rows / total_rows) * 100))

 

                            # Early exit if all keywords are found

                            if all(found_keywords.values()):

                                return True

 

                    # Check if all keywords are found after loop

                    return all(found_keywords.values())

 

                except Exception as e:

                    return None

 

            # For .pptx files

            elif file_extension == '.pptx':

                try:

                    prs = Presentation(file_path)

                    total_slides = len(prs.slides)

 

                    for i, slide in enumerate(prs.slides):

                        for shape in slide.shapes:

                            if hasattr(shape, "text"):

                                text = shape.text

                                for keyword in keywords:

                                    if keyword.lower() in text.lower():

                                        found_keywords[keyword.lower()] = True

                                    if keyword.lower() == '((regex))':

                                        for word in text.split():

                                            # Match the word against the compiled pattern

                                            if bool(pattern.fullmatch(word)):

                                                found_keywords['((regex))'] = True

                                                break

                                   

 

                        # Emit progress after processing each slide

                        self.progress_within_file.emit(int(((i + 1) / total_slides) * 100))

 

                    # Check if all keywords are found

                    return all(found_keywords.values())

 

                except Exception as e:

                    return None

 

            # For .msg files

            elif file_extension == '.msg':

                try:

                    msg = extract_msg.Message(file_path)

                    msg_content = msg.body.lower()

                    total_lines = len(msg_content.splitlines())

 

                    for i, line in enumerate(msg_content.splitlines()):

                        for keyword in keywords:

                            if keyword.lower() in line:

                                found_keywords[keyword.lower()] = True

                            if keyword.lower() == '((regex))':

                                for word in line.split():

                                    # Match the word against the compiled pattern

                                    if bool(pattern.fullmatch(word)):

                                        found_keywords['((regex))'] = True

                                        break

 

                        # Emit progress after processing each line

                        self.progress_within_file.emit(int(((i + 1) / total_lines) * 100))

 

                    # Check if all keywords are found

                    return all(found_keywords.values())

 

                except Exception as e:

                    return None

 

            # For .txt files

            elif file_extension == '.txt':

                try:

                    with open(file_path, 'r', encoding='utf-8') as file:

                        lines = file.readlines()

                        total_lines = len(lines)

 

                        for i, line in enumerate(lines):

                            for keyword in keywords:

                                if keyword.lower() in line.lower():

                                    found_keywords[keyword.lower()] = True

                                if keyword.lower() == '((regex))':

                                    for word in line.split():

                                        # Match the word against the compiled pattern

                                        if bool(pattern.fullmatch(word)):

                                            found_keywords['((regex))'] = True

                                            break

 

                            # Emit progress after processing each line

                            self.progress_within_file.emit(int(((i + 1) / total_lines) * 100))

 

                    # Check if all keywords are found

                    return all(found_keywords.values())

 

                except Exception as e:

                    return None

 

            # For .pdf files

            elif file_extension == '.pdf':

                try:

                    document = fitz.open(file_path)

                    total_pages = len(document)

 

                    for page_num in range(total_pages):

                        page = document.load_page(page_num)

 

                        for keyword in keywords:

                            if page.search_for(keyword):

                                found_keywords[keyword.lower()] = True

                            if keyword == '((regex))':

                                words = page.get_text("words")

                                # Extract all words on the page

                                for word_data in words:

                                    if bool(pattern.fullmatch(word_data[4])):

                                        found_keywords['((regex))'] = True

                                        break

 

                        # Emit progress after processing each page

                        self.progress_within_file.emit(int(((page_num + 1) / total_pages) * 100))

 

                    document.close()

 

                    # Check if all keywords are found

                    return all(found_keywords.values())

 

                except Exception as e:

                    return None

 

            # For .docx files

            elif file_extension == '.docx':

                try:

                    doc = Document(file_path)

 

                    for paragraph in doc.paragraphs:

                        for keyword in keywords:

                            if keyword.lower() in paragraph.text.lower():

                                found_keywords[keyword.lower()] = True

                            if keyword.lower() == '((regex))':

                                for word in paragraph.text.split():

                                    # Match the word against the compiled pattern

                                    if bool(pattern.fullmatch(word)):

                                        found_keywords['((regex))'] = True

                                        break

 

                    # Check if all keywords are found

                    return all(found_keywords.values())

 

                except Exception as e:

                    return None

 

        else:

            return False

 

        # Default return if no matches

        return False

 

 

class VerticalLabel(QLabel):

    def __init__(self, text, parent=None):

        super().__init__(text, parent)

 

    def paintEvent(self, event):

        painter = QPainter(self)

        painter.setPen(QColor(85, 85, 85))

        painter.setFont(QFont("Arial", 7))

 

        painter.translate(self.width()-3, self.height()-14)

        painter.rotate(-90)      

 

        painter.drawText(0, 0, self.text())

        painter.end()

 

    def sizeHint(self):

        return QSize(8, 1100)

 

 

class PDFViewer(QWidget):

    def __init__(self, parent=None):

        super().__init__(parent)

        self.zoom_level = 1.0

        self.current_page = 0

        self.doc = None

        self.initUI()

 

    def initUI(self):

        self.layout = QVBoxLayout()

 

        #--Credit label DAO-----------------------------------------------------------------------------#

 

        self.credit_label = QPushButton('Conceptualized and Developed by DAO India', self)

        self.credit_label.setStyleSheet("font-weight: bold")

        self.credit_label.setFont(QFont('Arial', 7))

        self.layout.addWidget(self.credit_label)

 

        #-----------------------------------------------------------------------------------------------#

 

        #--Scroll Area----------------------------------------------------------------------------------#

 

        self.scrollArea = QScrollArea()

        self.scrollWidget = QWidget()

        self.scrollLayout = QVBoxLayout()

        self.scrollLayout.setAlignment(Qt.AlignCenter)

        self.scrollWidget.setLayout(self.scrollLayout)

        self.scrollArea.setWidget(self.scrollWidget)

        self.scrollArea.setWidgetResizable(True)

        self.scrollArea.verticalScrollBar().valueChanged.connect(self.on_scroll)

        self.layout.addWidget(self.scrollArea)

 

        #----------------------------------------------------------------------------------------------#


 

        #--Bottom Functions----------------------------------------------------------------------------#

 

        self.button_layout = QHBoxLayout()

        # self.button_layout.addStretch()

 


        self.file_name_label = QLabel(self)

        self.file_name_label.setFixedSize(QSize(315, 18))

        self.button_layout.addWidget(self.file_name_label, alignment=Qt.AlignCenter | Qt.AlignLeft)


 

        self.page_label = QLabel(self)

        self.page_label.setFixedSize(QSize(150, 30))

        self.page_label.setAlignment(Qt.AlignRight | Qt.AlignCenter)

        self.button_layout.addWidget(self.page_label)


 

        self.zoom_out_button = QPushButton('-', self)

        self.zoom_out_button.setFixedSize(QSize(30, 30))

        self.zoom_out_button.clicked.connect(self.zoom_out)

        self.button_layout.addWidget(self.zoom_out_button)


 

        self.zoom_in_button = QPushButton('+', self)

        self.zoom_in_button.setFixedSize(QSize(30, 30))

        self.zoom_in_button.clicked.connect(self.zoom_in)

        self.button_layout.addWidget(self.zoom_in_button)


 

        self.layout.addLayout(self.button_layout)

 

        #----------------------------------------------------------------------------------------------#

 

        self.setLayout(self.layout)

        self.scrollArea.viewport().installEventFilter(self)

 


    def eventFilter(self, source, event):

        return super(PDFViewer, self).eventFilter(source, event)


 

    def clear_layout(self, layout):

        if layout is not None:

            while layout.count():

                child = layout.takeAt(0)

                if child.widget():

                    child.widget().deleteLater()


 

    def display_pdf(self, pdf_path, page_num=0):

        self.pdf_path = pdf_path

        self.file_name_label.setText(os.path.basename(pdf_path).split('_highlighted')[0]+'.pdf')

 

        if hasattr(self, 'pdf_path'):

            self.doc = fitz.open(self.pdf_path)

            self.clear_layout(self.scrollLayout)


 

            self.page_labels = []

 

            for p_num in range(len(self.doc)):

                self.add_page_to_layout(p_num)

            self.current_page = page_num

            self.scroll_to_page(page_num)

            self.update_page_label()

 


    def add_page_to_layout(self, page_num):

        page = self.doc.load_page(page_num)

        zoom_matrix = fitz.Matrix(self.zoom_level, self.zoom_level)

        pix = page.get_pixmap(matrix=zoom_matrix)

        image = QImage(pix.samples, pix.width, pix.height, pix.stride, QImage.Format_RGB888)

        pixmap = QPixmap.fromImage(image)

        label = QLabel()

        label.setPixmap(pixmap)

 

        if page_num < len(self.page_labels):

            self.scrollLayout.replaceWidget(self.page_labels[page_num], label)

            self.page_labels[page_num].deleteLater()

            self.page_labels[page_num] = label

        else:

            self.scrollLayout.addWidget(label)

            self.page_labels.append(label)

 


    def scroll_to_page(self, page_num):

        page_height = self.scrollArea.widget().height() / len(self.doc)

        self.scrollArea.verticalScrollBar().setValue(int(page_height * page_num))

        self.current_page = page_num

        self.update_page_label()


 

    def zoom_in(self):

        try:

            self.zoom_level += 0.1

            self.update_zoom()

        except:

            pass

 


    def zoom_out(self):

        try:

            if self.zoom_level > 0.1:

                self.zoom_level -= 0.1

            self.update_zoom()

        except:

            pass

 


    def update_zoom(self):

        for page_num in range(len(self.page_labels)):

            self.add_page_to_layout(page_num)

        self.update_page_label()

 


    def update_page_label(self):

        self.page_label.setText(f"Page {self.current_page+1} of {len(self.doc)}")

 


    def on_scroll(self):

        scroll_value = self.scrollArea.verticalScrollBar().value()

 

        try:

            page_height = self.scrollArea.widget().height() / len(self.doc)

            self.current_page = int(scroll_value / page_height)

            self.update_page_label()

        except:

            pass


 

    def close_pdf(self):

        if self.doc:

            self.doc.close()

            self.doc = None

            self.clear_layout(self.scrollLayout)


 

class PDFHighlighter(QMainWindow):

 

    def __init__(self):

        super().__init__()

        self.initUI()


 

    def initUI(self):

 

        self.setWindowTitle('Smart Key Search')

        self.setGeometry(100, 100, 1200, 900)

 

        self.setWindowIcon(QApplication.style().standardIcon(QApplication.style().SP_FileDialogContentsView))

        self.layout = QHBoxLayout()

        self.splitter = QSplitter(Qt.Horizontal)


 

        #--Credits Section----------------------------------------------------------------------------------#

 

        self.credit_label = VerticalLabel("For any queries contact | aaditya.shriram.dautkhane@hsbc.co.in | nilay.patel@hsbc.co.in | DAO - India", self)

        self.credit_label.setFixedSize(QSize(12, 800))

        self.layout.addWidget(self.credit_label, alignment=Qt.AlignRight | Qt.AlignBottom)

 

        #---------------------------------------------------------------------------------------------------

 

 

        #--Left Section-------------------------------------------------------------------------------------#

 

        self.left_panel = QWidget()

        self.left_layout = QVBoxLayout()

 

        #--Extension Selector------------------------------------------------------------------------------#

 

        self.extensionGroupBox = QGroupBox(' Select Extension ')

        self.extensionGroupBoxLayout = QVBoxLayout()

 

        self.extension_selector_row_1 = QHBoxLayout()

        self.extension_selector_row_2 = QHBoxLayout()

 

        self.pdf_button = QPushButton("PDF", self)

        self.docx_button = QPushButton("DOC", self)

        self.txt_button = QPushButton("TXT", self)

        self.xlsx_button = QPushButton("EXCEL", self)

        self.ppt_button = QPushButton("PPT", self)

        self.msg_button = QPushButton("OUTLOOK", self)


 

        # Add buttons to the layout

        self.extension_selector_row_1.addWidget(self.pdf_button)

        self.extension_selector_row_1.addWidget(self.docx_button)

        self.extension_selector_row_1.addWidget(self.txt_button)

        self.extension_selector_row_2.addWidget(self.xlsx_button)

        self.extension_selector_row_2.addWidget(self.ppt_button)

        self.extension_selector_row_2.addWidget(self.msg_button)


 

        # Connect the buttons to toggle functions

        self.pdf_button.clicked.connect(lambda: self.toggle_button(self.pdf_button))

        self.docx_button.clicked.connect(lambda: self.toggle_button(self.docx_button))

        self.txt_button.clicked.connect(lambda: self.toggle_button(self.txt_button))

        self.xlsx_button.clicked.connect(lambda: self.toggle_button(self.xlsx_button))

        self.ppt_button.clicked.connect(lambda: self.toggle_button(self.ppt_button))

        self.msg_button.clicked.connect(lambda: self.toggle_button(self.msg_button))


        # Initialize button states (selected)

        self.pdf_button.setProperty('selected', True)

        self.docx_button.setProperty('selected', True)

        self.txt_button.setProperty('selected', True)

        self.xlsx_button.setProperty('selected', True)

        self.ppt_button.setProperty('selected', True)

        self.msg_button.setProperty('selected', True)

 

        self.pdf_button.setStyleSheet("QPushButton { background-color: #D32F2F; color: white }")

        self.docx_button.setStyleSheet("QPushButton { background-color: #D32F2F; color: white }")

        self.txt_button.setStyleSheet("QPushButton { background-color: #D32F2F; color: white }")

        self.xlsx_button.setStyleSheet("QPushButton { background-color: #D32F2F; color: white }")

        self.ppt_button.setStyleSheet("QPushButton { background-color: #D32F2F; color: white }")

        self.msg_button.setStyleSheet("QPushButton { background-color: #D32F2F; color: white }")

 

        self.extensionGroupBoxLayout.addLayout(self.extension_selector_row_1)

        self.extensionGroupBoxLayout.addLayout(self.extension_selector_row_2)

 

        self.extensionGroupBox.setLayout(self.extensionGroupBoxLayout)

 

        self.left_layout.addWidget(self.extensionGroupBox)

 


        #--Search Section-----------------------------------------------------------------------------------#

 

        self.searchGroupBox = QGroupBox(' Search Folder ')

        self.searchGroupBoxLayout = QVBoxLayout()

 

        self.openFolderButton = QPushButton('Open Folder', self)

        self.openFolderButton.clicked.connect(self.open_folder)

        self.searchGroupBoxLayout.addWidget(self.openFolderButton)

 

        # Add the new QLabel to display the selected folder path

        self.folderPathLabel = QLabel('No folder selected', self)

        self.folderPathLabel.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Fixed)  # Expands horizontally

        self.folderPathLabel.setWordWrap(True)  # Allows the label to wrap text if the path is long

        self.folderPathLabel.setStyleSheet("background-color: #333333; border: 1px solid #666; border-radius: 5px; padding: 5px; color: #DDDDDD;")

 

        # Add the folderPathLabel to the layout below the "Open Folder" button

        self.searchGroupBoxLayout.addWidget(self.folderPathLabel)

 

        self.keywordInput = QTextEdit(self)

        self.keywordInput.setPlaceholderText('Enter keywords separated by commas')

        self.searchGroupBoxLayout.addWidget(self.keywordInput)

 

        self.regexInput = QTextEdit(self)

        self.regexInput.setPlaceholderText('Enter Regex')

        self.regexInput.setFixedHeight(38)

        self.searchGroupBoxLayout.addWidget(self.regexInput)

 

        self.file_progress_bar = QProgressBar(self)

        self.file_progress_bar.setVisible(False)  # Hide initially

        self.searchGroupBoxLayout.addWidget(self.file_progress_bar)

 

        self.total_progress_label = QLabel(self)

        self.total_progress_label.setVisible(False)  # Hide initially

        self.left_layout.addWidget(self.total_progress_label)

 

        self.searchButton = QPushButton('Search', self)

        self.searchButton.clicked.connect(self.search_keywords)

        self.searchGroupBoxLayout.addWidget(self.searchButton)

 

        self.searchGroupBox.setLayout(self.searchGroupBoxLayout)

        self.left_layout.addWidget(self.searchGroupBox)

 


 

        #--File list Section--------------------------------------------------------------------------------#

 

        self.fileListGroupBox = QGroupBox(' Results ')

        self.fileListGroupBoxLayout = QVBoxLayout()

 

        self.fileList = QListWidget(self)

        self.fileList.itemClicked.connect(self.open_selected_pdf)

        self.fileListGroupBoxLayout.addWidget(self.fileList)      

 

        self.fileListGroupBox.setLayout(self.fileListGroupBoxLayout)

 

        self.left_layout.addWidget(self.fileListGroupBox)

 

        #---------------------------------------------------------------------------------------------------

 

        # Label to show currently searched file

        self.current_file_label = QLabel(self)

        self.current_file_label.setAlignment(Qt.AlignLeft | Qt.AlignBottom)

        self.current_file_label.setFixedSize(QSize(300, 20))  # Adjust size as needed

        self.left_layout.addWidget(self.current_file_label)

 

        self.left_panel.setLayout(self.left_layout)

 

        #---------------------------------------------------------------------------------------------------

 

 

        #--Right Section------------------------------------------------------------------------------------#

 

        self.right_panel = QWidget()

        self.right_layout = QVBoxLayout()

 

        #--Toggle Mode-------------------------------------------------------------------------------------#

 

        self.mode_toggle_button = QPushButton('Toggle Dark/Light', self)

        self.mode_toggle_button.clicked.connect(self.toggle_mode)

        self.right_layout.addWidget(self.mode_toggle_button)

       

 

        #--Keyword List-------------------------------------------------------------------------------------#

 

        self.keyWordListGroupBox = QGroupBox(' Keyword List ')

        self.keyWordListGroupBoxLayout = QVBoxLayout()

 

        self.keywordList = QListWidget(self)

        self.keyWordListGroupBoxLayout.addWidget(self.keywordList)

        self.keywordList.itemClicked.connect(self.on_keyword_selected)

 

        self.up_down_button_layout = QHBoxLayout()

 

        self.upButton = QPushButton('Up', self)

        self.upButton.clicked.connect(self.scroll_keyword_up)

        self.up_down_button_layout.addWidget(self.upButton)

 

        self.downButton = QPushButton('Down', self)

        self.downButton.clicked.connect(self.scroll_keyword_down)

        self.up_down_button_layout.addWidget(self.downButton)

 

        self.keyWordListGroupBoxLayout.addLayout(self.up_down_button_layout)

 

        self.keyWordListGroupBox.setLayout(self.keyWordListGroupBoxLayout)

 

        self.right_layout.addWidget(self.keyWordListGroupBox)

 

        #---------------------------------------------------------------------------------------------------

 

 

        #--Page Line List----------------------------------------------------------------------------------#

 

        self.pageLineListGroupBox = QGroupBox(' Position Details ')

        self.pageLineListGroupBoxLayout = QVBoxLayout()

 

        self.page_line_list = QListWidget(self)

        self.page_line_list.itemClicked.connect(self.navigate_to_page_line)

        self.pageLineListGroupBoxLayout.addWidget(self.page_line_list)

 

        self.pageLineListGroupBox.setLayout(self.pageLineListGroupBoxLayout)

 

        self.right_layout.addWidget(self.pageLineListGroupBox)

 

        #---------------------------------------------------------------------------------------------------

 

        self.right_panel.setLayout(self.right_layout)

 

        #---------------------------------------------------------------------------------------------------

 

        self.pdf_viewer = PDFViewer()

 

        self.splitter.addWidget(self.left_panel)

        self.splitter.addWidget(self.pdf_viewer)

        self.splitter.addWidget(self.right_panel)

 

        # self.splitter.setStretchFactor(0, 1)

        # self.splitter.setStretchFactor(1, 3)

        # self.splitter.setStretchFactor(2, 1)

 

        self.splitter.setSizes([300, 800, 300])

        self.layout.addWidget(self.splitter)

 

        container = QWidget()

        container.setLayout(self.layout)

 

        self.setCentralWidget(container)

 

        self.temp_files = []

        self.page_positions = {}

        self.keyword_positions = {}

        self.keyword_list = []

        self.regex_string = ''

 

        self.apply_stylesheet()

 

    def toggle_button(self, button):

        # Check if the button is already selected

        if button.property('selected'):

            # Deselect the button, set normal color

            button.setStyleSheet("background-color: none")

            button.setProperty('selected', False)

        else:

            # Select the button, set to red

            button.setStyleSheet("QPushButton { background-color: #D32F2F; color: white }")

            button.setProperty('selected', True)


 

    def toggle_mode(self):

        if hasattr(self, 'dark_mode_enabled') and self.dark_mode_enabled:

            self.setStyleSheet(self.light_mode_stylesheet)

            self.folderPathLabel.setStyleSheet("background-color: #F0F0F0; border: 1px solid #999; border-radius: 5px; padding: 5px;")

            self.dark_mode_enabled = False

        else:

            self.setStyleSheet(self.dark_mode_stylesheet)

            self.folderPathLabel.setStyleSheet("background-color: #333333; border: 1px solid #666; border-radius: 5px; padding: 5px; color: #DDDDDD;")

            self.dark_mode_enabled = True


 

    def apply_stylesheet(self):

        light_mode_stylesheet = """

            QWidget {

                background-color: #F0F0F0;

                color: #333;

                font-family: Arial, sans-serif;

            }

            QPushButton {

                background-color: #E0E0E0;

                border: 1px solid #999;

                padding: 5px;

                border-radius: 5px;

            }

            QPushButton#searchButton {

                background-color: #D32F2F;

                color: white;

            }

            QListWidget::item:selected {

                background-color: #D32F2F;

                color: white;

                border: 1px solid #D32F2F;

            }

            QPushButton:hover {

                background-color: #C0C0C0;

            }

            QTextEdit, QListWidget {

                background-color: #FFFFFF;

                border: 1px solid #999;

                padding: 5px;

                border-radius: 5px;

            }

            QGroupBox {

                border: 1px solid #999;

                padding-top: 10px;

                border-radius: 5px;

            }

            QSplitter::handle {

                background-color: #858484;

                width: 1px;

                height: 1px;

                border-radius: 1px;

            }

            QProgressBar {

                border: 1px solid #999;

                border-radius: 5px;

                background-color: #E0E0E0;

                text-align: center;

                height: 10px;  /* Adjust the height */

            }

            QProgressBar::chunk {

                background-color: #4CAF50;  /* Green progress chunk */

                width: 10px;  /* Chunk width */

            }

        """


        dark_mode_stylesheet = """

            QWidget {

                background-color: #2C2C2C;

                color: #DDDDDD;

                font-family: Arial, sans-serif;

            }

            QPushButton {

                background-color: #444444;

                border: 1px solid #666;

                padding: 5px;

                border-radius: 5px;

            }

            QPushButton#searchButton {

                background-color: #8B0000;

                color: white;

            }

            QListWidget::item:selected {

                background-color: #8B0000;

                color: white;

                border: 1px solid #8B0000;

            }

            QPushButton:hover {

                background-color: #555555;

            }

            QTextEdit, QListWidget {

                background-color: #333333;

                border: 1px solid #666;

                padding: 5px;

                border-radius: 5px;

            }

            QGroupBox {

                border: 1px solid #666;

                padding-top: 10px;

                border-radius: 5px;

            }

            QSplitter::handle {

                background-color: #444444;

                width: 1px;

                height: 1px;

                border-radius: 1px;

            }

            QProgressBar {

                border: 1px solid #666;

                border-radius: 5px;

                background-color: #2E2E2E;

                text-align: center;

                height: 10px;  /* Adjust the height */

            }

            QProgressBar::chunk {

                background-color: #2196F3;  /* Blue progress chunk */

                width: 10px;  /* Chunk width */

            }

        """

 

        self.light_mode_stylesheet = light_mode_stylesheet

        self.dark_mode_stylesheet = dark_mode_stylesheet

        self.dark_mode_enabled = True  # Start in light mode

 

        # Apply light mode initially

        self.setStyleSheet(self.dark_mode_stylesheet)

 

 

    def open_folder(self):

        options = QFileDialog.Options()

        self.folder_path = QFileDialog.getExistingDirectory(self, 'Open Folder', '', options=options)

 

        if self.folder_path:

            self.folderPathLabel.setText(self.folder_path)

        else:

            self.folderPathLabel.setText('No folder selected')

 

        if self.folder_path:

            self.keywordList.clear()

            self.page_line_list.clear()

            self.fileList.clear()

            self.pdf_viewer.close_pdf()

 

            for temp_file in self.temp_files:

                try:

                    os.remove(temp_file)

                    self.temp_files = []

                except:

                    pass


 

    def convert_txt_to_pdf(self, txt_path):

        pdf_path = txt_path.rsplit('.', 1)[0] + '.pdf'      

 

        with open(txt_path, 'r', encoding='utf-8') as file:

            content = file.read()

 

        pdf = SimpleDocTemplate(pdf_path, pagesize=letter)

        styles = getSampleStyleSheet()

        story = [Paragraph(content.replace('\n', '<br/>'), styles['Normal'])]

        pdf.build(story)

 

        self.temp_files.append(pdf_path)

 

        return pdf_path

 

 

    def search_keywords(self):

        # Clear previous search data

        self.keywordList.clear()

        self.page_line_list.clear()

        self.pdf_viewer.close_pdf()

 

        # Remove temporary files

        for temp_file in self.temp_files:

            try:

                os.remove(temp_file)

            except:

                pass

 

        self.temp_files = []

        self.keyword_list = []

        self.regex_string = ''

 

        # Check if a folder is selected

        if hasattr(self, 'folder_path'):

            keywords = self.keywordInput.toPlainText().split(',')

            self.regex_string = self.regexInput.toPlainText()

 

            self.keyword_list = [k.strip() for k in keywords]

 

            if self.regex_string != '':

                self.keyword_list.append('((regex))')

 

            self.fileList.clear()

 

            # if '' not in keywords:

            # Get selected file formats based on button states

            file_extensions = []

            if self.pdf_button.property('selected'):

                file_extensions.append('.pdf')

            if self.docx_button.property('selected'):

                file_extensions.append('.docx')

            if self.txt_button.property('selected'):

                file_extensions.append('.txt')

            if self.xlsx_button.property('selected'):

                file_extensions.append('.xlsx')

            if self.ppt_button.property('selected'):

                file_extensions.append('.pptx')

            if self.msg_button.property('selected'):

                file_extensions.append('.msg')

 

            # Start the search in a background thread

            self.file_progress_bar.setVisible(True)

            self.total_progress_label.setVisible(True)

            self.current_file_label.setVisible(True)

 

            # Initialize and start the SearchThread

            self.search_thread = SearchThread(self.folder_path, self.keyword_list, self.regex_string, file_extensions)

            self.search_thread.progress_within_file.connect(self.update_file_progress_bar)

            self.search_thread.update_current_file.connect(self.update_current_file_label)

            self.search_thread.file_found.connect(self.add_file_to_list)  # Connect the file_found signal

            self.search_thread.search_complete.connect(self.on_search_complete)

            self.search_thread.files_searched.connect(self.update_files_searched_label)  # Connect new signal

 

            # Start the thread

            self.search_thread.start()

 

            # else:

            #     self.current_file_label.setText("Empty Keyword Error! Check After Comma.")

 

        else:

            self.current_file_label.setText("Folder Path NOT Selected.")

 

 

    def update_file_progress_bar(self, progress):

        self.file_progress_bar.setValue(progress)

 

 

    def update_files_searched_label(self, current_file_index, total_files):

        self.total_progress_label.setText(f"Files searched: {current_file_index}/{total_files}")

 

 

    def update_current_file_label(self, file_path):

        self.current_file_label.setText(f"Searching: {os.path.basename(file_path)}")

 

 

    def add_file_to_list(self, file_path):

        self.fileList.addItem(file_path)  # Add the file with matching keywords to the file list

 

 

    def on_search_complete(self):

        self.file_progress_bar.setVisible(False)

        self.total_progress_label.setVisible(False)

        if self.fileList.count():

            self.current_file_label.setText("Search Complete.")

        else:

            self.current_file_label.setText("All Keywords NOT Found.")

 


    def open_selected_pdf(self, item):

        self.page_line_list.clear()

 

        # Use the full path of the selected file from the item

        full_file_path = item.text()

        file_name, file_extension = os.path.splitext(full_file_path)

        file_extension = file_extension.lower()

        existing_highlighted_pdf_path = f"{file_name}_highlighted_{'_'.join(self.keyword_list)}.pdf"

 

        if existing_highlighted_pdf_path in self.temp_files:

            self.on_highlight_complete(existing_highlighted_pdf_path)

            return None

 

        if file_extension == '.txt':

            full_file_path = self.convert_txt_to_pdf(full_file_path)

 

        # Start a new thread to highlight the selected file

        self.highlight_thread = HighlightThread(full_file_path, self.keyword_list, self.regex_string, file_extension)

 

        # Connect signals to handle completion and errors

        self.highlight_thread.highlight_complete.connect(self.on_highlight_complete)

        self.highlight_thread.error.connect(self.on_highlight_error)

 

        # Start the thread (this runs the highlighting in the background)

        self.highlight_thread.start()

 

    def on_highlight_complete(self, highlighted_file_path):

        self.temp_files.append(highlighted_file_path)

        current_page = self.page_positions.get(highlighted_file_path, 0)

        self.pdf_viewer.display_pdf(highlighted_file_path, current_page)

        self.keyword_positions = self.get_keyword_positions(highlighted_file_path, self.keyword_list)

        self.populate_keyword_list(self.keyword_list)

 

    def on_highlight_error(self, error_message):

        print(f"Error during highlighting: {error_message}")

 

    def get_keyword_positions(self, pdf_path, keywords):

        document = fitz.open(pdf_path)

        keyword_positions = {keyword: [] for keyword in keywords}

 

        for page_num in range(len(document)):

            page = document.load_page(page_num)

            blocks = page.get_text("dict")["blocks"]

 

            lines = []

            line_number = 1

 

            for block in blocks:

                for line in block.get("lines", []):

                    line_bbox = line["bbox"]

                    line_y0 = line_bbox[1]

                    line_y1 = line_bbox[3]

                    line_text = " ".join([span["text"] for span in line["spans"] if span["text"].strip() != ""])

 

                    if line_text.strip():

                        lines.append((line_number, line_y0, line_y1, line_text))

                        line_number += 1

 

            for keyword in keywords:

                if keyword == '((regex))':

                    words = page.get_text("words")  # Extract all words on the page

                    for word_data in words:

                        if bool(re.compile(self.regex_string).fullmatch(word_data[4])):

                            # inst = fitz.Rect(word_data[0], word_data[1], word_data[2], word_data[3])

                            y0 = word_data[1]

                            y1 = word_data[3]

                            line_text = None

 

                            for line_num, line_y0, line_y1, text in lines:

                                if abs(line_y0 - y0) < 2 and abs(line_y1 - y1) < 2:

                                    line_text = text

                                    break

 

                            if line_text:

                                snippet = line_text[:20] + "..."

                                keyword_positions[keyword].append((page_num, line_num, snippet))

                elif keyword != '':

                    text_instances = page.search_for(keyword)

                    for inst in text_instances:

                        y0 = inst.y0

                        y1 = inst.y1

                        line_text = None

 

                        for line_num, line_y0, line_y1, text in lines:

                            if abs(line_y0 - y0) < 2 and abs(line_y1 - y1) < 2:

                                line_text = text

                                break

 

                        if line_text:

                            snippet = line_text[:20] + "..."

                            keyword_positions[keyword].append((page_num, line_num, snippet))

 

        document.close()

 

        return keyword_positions


 

    def populate_keyword_list(self, keywords):

        self.keywordList.clear()

        colors = [QColor("magenta"), QColor("red"), QColor("blue"), QColor("pink"), QColor("orange"), QColor("green"), QColor("yellow"), QColor("cyan")]

 

        for i, keyword in enumerate(keywords):

            if keyword == '((regex))':

                item = QListWidgetItem(keyword)

                item.setBackground(colors[len(keywords)-1 % len(colors)])

                self.keywordList.addItem(item)

            elif keyword != '':

                item = QListWidgetItem(keyword)

                item.setBackground(colors[i % len(colors)])

                self.keywordList.addItem(item)

           

 

    def on_keyword_selected(self, item):

        selected_keyword = item.text()

        self.populate_page_line_list(selected_keyword)

 

    def scroll_keyword_up(self):

        selected_items = self.keywordList.selectedItems()

        if not selected_items:

            return

 

        selected_keyword = selected_items[0].text()

 

        if selected_keyword in self.keyword_positions:

            current_page = self.pdf_viewer.current_page

            pages = [pos[0] for pos in self.keyword_positions[selected_keyword]]

            previous_pages = [p for p in pages if p < current_page]

 

            if previous_pages:

                self.pdf_viewer.scroll_to_page(previous_pages[-1])

            else:

                self.pdf_viewer.scroll_to_page(pages[-1])

 

            self.populate_page_line_list(selected_keyword)

 

    def scroll_keyword_down(self):

        selected_items = self.keywordList.selectedItems()

 

        if not selected_items:

            return

 

        selected_keyword = selected_items[0].text()

        if selected_keyword in self.keyword_positions:

            current_page = self.pdf_viewer.current_page

            pages = [pos[0] for pos in self.keyword_positions[selected_keyword]]

            next_pages = [p for p in pages if p > current_page]

            if next_pages:

                self.pdf_viewer.scroll_to_page(next_pages[0])

            else:

                self.pdf_viewer.scroll_to_page(pages[0])

 

            self.populate_page_line_list(selected_keyword)

 


    def populate_page_line_list(self, keyword):

        self.page_line_list.clear()

        if keyword in self.keyword_positions:

            for page_num, y_pos, snippet in self.keyword_positions[keyword]:

                item = QListWidgetItem(f"Page {page_num + 1}: Line {int(y_pos)} | {snippet}")

                item.setData(Qt.UserRole, (page_num, y_pos))

                self.page_line_list.addItem(item)

 


    def navigate_to_page_line(self, item):

        page_num, y_pos = item.data(Qt.UserRole)

        self.pdf_viewer.scroll_to_page(page_num)

        scroll_value = self.pdf_viewer.scrollArea.verticalScrollBar().value() + y_pos + 2

        self.pdf_viewer.scrollArea.verticalScrollBar().setValue(scroll_value)

 


    def closeEvent(self, event):

        self.pdf_viewer.close_pdf()

        for temp_file in self.temp_files:

            try:

                os.remove(temp_file)

            except:

                pass

        event.accept()

 


if __name__ == '__main__':

    app = QApplication(sys.argv)

    ex = PDFHighlighter()

    ex.show()

    sys.exit(app.exec_())
